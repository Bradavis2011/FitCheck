#!/usr/bin/env python3
"""
Generate the "Or This?" Seed Stage Pitch Deck as a PowerPoint file.
Design language reverse-engineered from OrThis_Seed_Pitch_Deck(2).pptx reference.
Designed for Y Combinator and similar accelerator applications.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand Colors ──
CORAL       = RGBColor(0xE8, 0x5D, 0x4C)   # #E85D4C — Decision Coral
BLACK       = RGBColor(0x1A, 0x1A, 0x1A)   # #1A1A1A — Clarity Black
CHARCOAL    = RGBColor(0x2D, 0x2D, 0x2D)   # #2D2D2D
GRAY        = RGBColor(0x9B, 0x9B, 0x9B)   # #9B9B9B — secondary / muted
DIVIDER     = RGBColor(0xE8, 0xE8, 0xE8)   # #E8E8E8 — thin rule
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)   # #FFFFFF
CREAM       = RGBColor(0xFB, 0xF7, 0xF4)   # #FBF7F4 — fallback placeholder

# ── Slide Dimensions (16:9) ──
SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ── Image paths (extracted from reference deck) ──
IMG_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "extracted_images")

def img(filename):
    path = os.path.join(IMG_DIR, filename)
    return path if os.path.exists(path) else None


# ──────────────────────────────────────────────────────────────────────────────
# Primitive helpers
# ──────────────────────────────────────────────────────────────────────────────

def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None,
              shape_type=MSO_SHAPE.RECTANGLE, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_picture(slide, path, left, top, width, height, placeholder=CREAM):
    if path and os.path.exists(path):
        return slide.shapes.add_picture(path, left, top, width, height)
    # Fallback: solid color rectangle
    return add_shape(slide, left, top, width, height, placeholder)


def add_text_box(slide, left, top, width, height, text,
                 font_size=16, font_color=CHARCOAL, bold=False, italic=False,
                 alignment=PP_ALIGN.LEFT, font_name="DM Sans", line_spacing=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    if line_spacing:
        p.line_spacing = Pt(font_size * line_spacing)
    return tb


def add_multiline_text(slide, left, top, width, height, lines,
                       default_size=16, default_color=CHARCOAL, default_bold=False,
                       alignment=PP_ALIGN.LEFT, font_name="DM Sans", line_spacing=1.4):
    """
    lines: list of str or dict.
    dict keys: text, size, color, bold, italic, font_name, space_after
    """
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

        if isinstance(line, str):
            p.text = line
            p.font.size = Pt(default_size)
            p.font.color.rgb = default_color
            p.font.bold = default_bold
            p.font.italic = False
            p.font.name = font_name
            p.space_after = Pt(4)
        else:
            p.text = line.get("text", "")
            sz = line.get("size", default_size)
            p.font.size = Pt(sz)
            p.font.color.rgb = line.get("color", default_color)
            p.font.bold = line.get("bold", default_bold)
            p.font.italic = line.get("italic", False)
            p.font.name = line.get("font_name", font_name)
            p.space_after = Pt(line.get("space_after", 4))

        p.alignment = alignment
        if line_spacing and (isinstance(line, str) or not line.get("no_spacing")):
            sz = default_size if isinstance(line, str) else line.get("size", default_size)
            p.line_spacing = Pt(sz * line_spacing)

    return tb


# ── Short editorial coral rule (placed below headline, not above section label) ──
def coral_rule(slide, left, top, width=Inches(1.5)):
    add_shape(slide, left, top, width, Inches(0.028), CORAL)


# ── Full-width thin gray divider (before footnote at slide bottom) ──
def gray_rule(slide, left=Inches(0.8), top=None, width=Inches(11.5)):
    add_shape(slide, left, top, width, Inches(0.014), DIVIDER)


# ── Card with outline (for FREE tier, white on white bg) ──
def add_outline_card(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = DIVIDER
    shape.line.width = Pt(1)
    return shape


# ── Internal thin divider inside a card ──
def card_divider(slide, left, top, width, color=DIVIDER):
    add_shape(slide, left, top, width, Inches(0.014), color)


# ──────────────────────────────────────────────────────────────────────────────
# Slide builders
# ──────────────────────────────────────────────────────────────────────────────

def build_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank = prs.slide_layouts[6]

    # ══════════════════════════════════════════════════════════════
    # SLIDE 1: TITLE  — white bg, left half image, right text
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)

    # Left half — fashion photo panel
    add_picture(slide, img("slide1_Picture 1.jpg"),
                Inches(0), Inches(0), Inches(5.5), Inches(7.5))

    # Right side — coral rule above the logo
    coral_rule(slide, left=Inches(6.5), top=Inches(2.0), width=Inches(2.0))

    # "Or " — DM Sans 72pt black
    add_text_box(slide, Inches(5.88), Inches(2.3), Inches(1.48), Inches(1.31),
                 "Or", font_size=72, font_color=BLACK, font_name="DM Sans")

    # "This?" — Playfair Display 72pt coral
    add_text_box(slide, Inches(7.36), Inches(2.3), Inches(2.59), Inches(1.31),
                 "This?", font_size=72, font_color=CORAL, font_name="Playfair Display")

    # Tagline
    add_text_box(slide, Inches(6.5), Inches(3.7), Inches(6.0), Inches(0.6),
                 "Confidence in every choice.",
                 font_size=24, font_color=CHARCOAL, font_name="Playfair Display")

    # One-liner
    add_multiline_text(slide, Inches(6.5), Inches(4.6), Inches(5.5), Inches(0.8), [
        {"text": "AI-powered outfit feedback in 30 seconds.", "size": 16, "color": GRAY},
        {"text": "Your honest friend in your pocket.",        "size": 16, "color": GRAY},
    ])

    # Stage line
    add_text_box(slide, Inches(6.5), Inches(6.0), Inches(5.0), Inches(0.4),
                 "Seed Stage  \u2022  2026",
                 font_size=13, font_color=GRAY, font_name="DM Sans")

    # ══════════════════════════════════════════════════════════════
    # SLIDE 2: THE PROBLEM  — white bg, left text, right inset photo
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)

    # Right inset photo
    add_picture(slide, img("slide2_Picture 6.jpg"),
                Inches(8.8), Inches(0.8), Inches(4.2), Inches(6.2))

    # Section label
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5.0), Inches(0.4),
                 "THE PROBLEM", font_size=11, font_color=CORAL, font_name="DM Sans")

    # Headline
    add_multiline_text(slide, Inches(0.8), Inches(1.2), Inches(7.5), Inches(1.4), [
        {"text": "Every morning, millions of women", "size": 32, "color": BLACK, "font_name": "Playfair Display"},
        {"text": "stand in front of the mirror and ask:", "size": 32, "color": BLACK, "font_name": "Playfair Display"},
        {"text": "\u201cDoes this actually look good?\u201d", "size": 32, "color": BLACK, "font_name": "Playfair Display"},
    ], line_spacing=1.2)

    # Thin divider between headline and bullets
    gray_rule(slide, left=Inches(0.8), top=Inches(3.0), width=Inches(1.5))

    # Bullet stats
    add_multiline_text(slide, Inches(0.8), Inches(3.3), Inches(5.5), Inches(2.2), [
        {"text": "72% of women say outfit indecision causes daily stress.", "size": 18, "color": CHARCOAL},
        {"text": "The average woman changes outfits 2\u20133 times before leaving.", "size": 18, "color": CHARCOAL},
        {"text": "No trusted, instant feedback exists at the moment of decision.", "size": 18, "color": CHARCOAL},
    ], line_spacing=1.5)

    # Pull quote
    add_text_box(slide, Inches(0.8), Inches(6.2), Inches(7.0), Inches(0.6),
                 "\u2014 \u201cI just want someone honest to tell me if this works.\u201d",
                 font_size=17, font_color=CHARCOAL, font_name="Playfair Display")

    # ══════════════════════════════════════════════════════════════
    # SLIDE 3: THE SOLUTION  — white bg, left half image, right steps
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)

    add_picture(slide, img("slide3_Picture 1.jpg"),
                Inches(0), Inches(0), Inches(5.5), Inches(7.5))

    add_text_box(slide, Inches(6.2), Inches(0.5), Inches(5.0), Inches(0.4),
                 "THE SOLUTION", font_size=11, font_color=CORAL, font_name="DM Sans")

    add_multiline_text(slide, Inches(6.2), Inches(1.2), Inches(6.5), Inches(1.0), [
        {"text": "Your honest friend", "size": 36, "color": BLACK, "font_name": "Playfair Display"},
        {"text": "in your pocket.",    "size": 36, "color": BLACK, "font_name": "Playfair Display"},
    ], line_spacing=1.15)

    coral_rule(slide, left=Inches(6.2), top=Inches(2.6))

    # 4-step process
    steps = [
        ("01", "Snap",         "Full-body photo or selfie."),
        ("02", "Context",      "Occasion, setting, vibe."),
        ("03", "Feedback",     "AI analysis in 30 seconds."),
        ("04", "Conversation", "Follow-ups like a friend."),
    ]
    for i, (num, title, desc) in enumerate(steps):
        y = Inches(3.0 + i * 1.0)
        add_text_box(slide, Inches(6.2), y, Inches(0.7), Inches(0.4),
                     num, font_size=14, font_color=DIVIDER, font_name="Playfair Display")
        add_text_box(slide, Inches(7.0), y, Inches(2.0), Inches(0.4),
                     title, font_size=18, font_color=BLACK, font_name="DM Sans")
        add_text_box(slide, Inches(9.2), y, Inches(3.5), Inches(0.4),
                     desc, font_size=15, font_color=CHARCOAL, font_name="DM Sans")

    gray_rule(slide, left=Inches(6.2), top=Inches(6.6), width=Inches(6.5))
    add_text_box(slide, Inches(6.2), Inches(6.8), Inches(6.5), Inches(0.5),
                 "Not inspiration. Not shopping. Just honest, instant, private feedback.",
                 font_size=15, font_color=CHARCOAL, font_name="Playfair Display")

    # ══════════════════════════════════════════════════════════════
    # SLIDE 4: THE AI  — dark bg, left half image, right bullets
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BLACK)

    add_picture(slide, img("slide4_Picture 1.jpg"),
                Inches(0), Inches(0), Inches(5.5), Inches(7.5))

    add_text_box(slide, Inches(6.2), Inches(0.5), Inches(5.0), Inches(0.4),
                 "THE AI", font_size=11, font_color=GRAY, font_name="DM Sans")

    add_text_box(slide, Inches(6.2), Inches(1.2), Inches(6.5), Inches(0.8),
                 "It tells you the truth.",
                 font_size=40, font_color=WHITE, font_name="Playfair Display")

    add_text_box(slide, Inches(6.2), Inches(2.1), Inches(6.0), Inches(0.4),
                 "A Vogue editor in your pocket \u2014 powered by vision AI.",
                 font_size=16, font_color=GRAY, font_name="DM Sans")

    coral_rule(slide, left=Inches(6.2), top=Inches(2.7))

    add_multiline_text(slide, Inches(6.2), Inches(3.1), Inches(6.0), Inches(2.5), [
        {"text": "\u2014  A score out of 10. No sugarcoating.",          "size": 16, "color": GRAY},
        {"text": "\u2014  What\u2019s working \u2014 specific, actionable.",   "size": 16, "color": GRAY, "space_after": 6},
        {"text": "\u2014  What to reconsider \u2014 honest, constructive.", "size": 16, "color": GRAY, "space_after": 6},
        {"text": "\u2014  Quick fixes you can change in 2 minutes.",      "size": 16, "color": GRAY, "space_after": 6},
        {"text": "\u2014  Keep asking until you\u2019re sure.",            "size": 16, "color": GRAY, "space_after": 12},
    ])

    add_text_box(slide, Inches(6.2), Inches(5.3), Inches(6.0), Inches(0.3),
                 "Built-in fashion expertise:",
                 font_size=14, font_color=WHITE, font_name="DM Sans")

    add_text_box(slide, Inches(6.2), Inches(5.7), Inches(6.5), Inches(0.35),
                 "Color theory  \u2022  Proportions  \u2022  Fit  \u2022  Dress codes  \u2022  Trends",
                 font_size=13, font_color=GRAY, font_name="DM Sans")

    add_text_box(slide, Inches(6.2), Inches(6.1), Inches(6.5), Inches(0.35),
                 "Automated 7-stage training pipeline  \u2022  Prompt v3.0",
                 font_size=13, font_color=GRAY, font_name="DM Sans")

    # ══════════════════════════════════════════════════════════════
    # SLIDE 5: MARKET OPPORTUNITY  — white bg, left dark cards, right image
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)

    add_picture(slide, img("slide5_Image 0.jpg"),
                Inches(8.8), Inches(0), Inches(4.53), Inches(7.5))

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5.0), Inches(0.4),
                 "MARKET OPPORTUNITY", font_size=11, font_color=CORAL, font_name="DM Sans")

    add_multiline_text(slide, Inches(0.8), Inches(0.9), Inches(7.8), Inches(0.8), [
        {"text": "Two markets. One product.", "size": 36, "color": BLACK, "font_name": "Playfair Display"},
    ], line_spacing=1.2)

    add_text_box(slide, Inches(0.8), Inches(1.8), Inches(7.8), Inches(0.5),
                 "Consumer subscriptions + B2B fashion intelligence licensing.",
                 font_size=17, font_color=CHARCOAL, font_name="DM Sans")

    coral_rule(slide, left=Inches(0.8), top=Inches(2.5))

    # TAM / SAM / SOM — dark cards
    markets = [
        ("TAM",  "$1.8T",  "Global apparel\nmarket",              Inches(0.8)),
        ("SAM",  "$5.5B",  "Online personal styling\n(+$15B by 2032, 15% CAGR)", Inches(3.4)),
        ("SOM",  "$200M",  "Year 5 ARR\n2M subs \u00d7 $100/yr + B2B data", Inches(6.0)),
    ]
    for label, value, desc, x in markets:
        add_shape(slide, x, Inches(2.9), Inches(2.4), Inches(2.6), BLACK)
        add_text_box(slide, x + Inches(0.2), Inches(3.1), Inches(2.0), Inches(0.3),
                     label, font_size=11, font_color=GRAY, font_name="DM Sans")
        add_text_box(slide, x + Inches(0.2), Inches(3.4), Inches(2.0), Inches(0.7),
                     value, font_size=40, font_color=WHITE, font_name="Playfair Display")
        add_text_box(slide, x + Inches(0.2), Inches(4.2), Inches(2.0), Inches(0.7),
                     desc, font_size=13, font_color=GRAY, font_name="DM Sans")

    gray_rule(slide, left=Inches(0.8), top=Inches(6.2))
    add_text_box(slide, Inches(0.8), Inches(6.4), Inches(7.8), Inches(0.7),
                 "Why now?  Vision AI crossed the quality threshold in 2023\u201324. "
                 "And every outfit we analyze becomes a permanent, structured data asset "
                 "that compounds in value \u2014 a moat no competitor can replicate without our user base.",
                 font_size=15, font_color=CHARCOAL, font_name="Playfair Display")

    # ══════════════════════════════════════════════════════════════
    # SLIDE 6: FASHION INTELLIGENCE  — dark bg, data platform story
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BLACK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11.0), Inches(0.4),
                 "FASHION INTELLIGENCE", font_size=11, font_color=GRAY, font_name="DM Sans")

    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.9),
                 "The consumer app is the data collection engine.",
                 font_size=40, font_color=WHITE, font_name="Playfair Display")

    add_text_box(slide, Inches(0.8), Inches(2.1), Inches(11.5), Inches(0.45),
                 "Every outfit check generates structured fashion intelligence no competitor "
                 "can replicate without the same user base.",
                 font_size=16, font_color=GRAY, font_name="DM Sans")

    coral_rule(slide, left=Inches(0.8), top=Inches(2.75))

    # Three-panel pipeline: CAPTURE → DATASET → B2B MARKET
    DARK_CARD = RGBColor(0x2D, 0x2D, 0x2D)

    # Panel 1 — WHAT WE CAPTURE
    add_shape(slide, Inches(0.8), Inches(3.0), Inches(3.6), Inches(2.9), DARK_CARD)
    add_text_box(slide, Inches(1.0), Inches(3.15), Inches(3.2), Inches(0.3),
                 "CAPTURED PER CHECK", font_size=11, font_color=CORAL, font_name="DM Sans")
    add_multiline_text(slide, Inches(1.0), Inches(3.55), Inches(3.2), Inches(2.2), [
        {"text": "\u2014  StyleDNA attributes",            "size": 14, "color": GRAY, "space_after": 5},
        {"text": "\u2014  Color harmonies & palette",      "size": 14, "color": GRAY, "space_after": 5},
        {"text": "\u2014  Occasion + setting patterns",    "size": 14, "color": GRAY, "space_after": 5},
        {"text": "\u2014  Garment categories & frequency", "size": 14, "color": GRAY, "space_after": 5},
        {"text": "\u2014  AI score + community signal",   "size": 14, "color": GRAY, "space_after": 5},
        {"text": "\u2014  Trend velocity signals",         "size": 14, "color": GRAY},
    ])

    # Arrow label between panels
    add_text_box(slide, Inches(4.5), Inches(4.3), Inches(0.5), Inches(0.4),
                 "\u2192", font_size=24, font_color=CORAL, font_name="DM Sans",
                 alignment=PP_ALIGN.CENTER)

    # Panel 2 — AT SCALE
    add_shape(slide, Inches(5.0), Inches(3.0), Inches(3.6), Inches(2.9), DARK_CARD)
    add_text_box(slide, Inches(5.2), Inches(3.15), Inches(3.2), Inches(0.3),
                 "AT SCALE", font_size=11, font_color=CORAL, font_name="DM Sans")
    add_multiline_text(slide, Inches(5.2), Inches(3.55), Inches(3.2), Inches(2.2), [
        {"text": "Real-time, bottom-up dataset",            "size": 14, "color": WHITE, "space_after": 8},
        {"text": "What real people actually wear \u2014",   "size": 13, "color": GRAY, "space_after": 4},
        {"text": "not what they say they wear.",            "size": 13, "color": GRAY, "space_after": 10},
        {"text": "AI-scored. Community-validated.",         "size": 13, "color": GRAY, "space_after": 4},
        {"text": "Queryable by demo, occasion, geo.",       "size": 13, "color": GRAY, "space_after": 10},
        {"text": "No competitor can replicate it",          "size": 13, "color": GRAY, "space_after": 4},
        {"text": "without our user base.",                  "size": 13, "color": GRAY},
    ])

    # Arrow label between panels
    add_text_box(slide, Inches(8.7), Inches(4.3), Inches(0.5), Inches(0.4),
                 "\u2192", font_size=24, font_color=CORAL, font_name="DM Sans",
                 alignment=PP_ALIGN.CENTER)

    # Panel 3 — B2B MARKET
    add_shape(slide, Inches(9.2), Inches(3.0), Inches(3.6), Inches(2.9), DARK_CARD)
    add_text_box(slide, Inches(9.4), Inches(3.15), Inches(3.2), Inches(0.3),
                 "B2B REVENUE", font_size=11, font_color=CORAL, font_name="DM Sans")
    add_multiline_text(slide, Inches(9.4), Inches(3.55), Inches(3.2), Inches(2.2), [
        {"text": "Trend reports",                   "size": 14, "color": WHITE, "space_after": 2},
        {"text": "$15\u201350K / quarter",           "size": 13, "color": GRAY, "space_after": 7},
        {"text": "Custom brand research",           "size": 14, "color": WHITE, "space_after": 2},
        {"text": "$25\u2013100K per study",          "size": 13, "color": GRAY, "space_after": 7},
        {"text": "Real-time trend API",             "size": 14, "color": WHITE, "space_after": 2},
        {"text": "$5\u201320K / month",              "size": 13, "color": GRAY, "space_after": 7},
        {"text": "White-label styling SDK",         "size": 14, "color": WHITE, "space_after": 2},
        {"text": "$10\u201350K / month",             "size": 13, "color": GRAY},
    ])

    # Bottom insight — WGSN comparable
    gray_rule(slide, left=Inches(0.8), top=Inches(6.2))
    add_text_box(slide, Inches(0.8), Inches(6.4), Inches(12.2), Inches(0.7),
                 "Comparable: WGSN earns ~$100M+ ARR selling fashion intelligence "
                 "built on runway predictions and surveys. "
                 "Ours is built on what real people actually wear.",
                 font_size=16, font_color=CORAL, font_name="Playfair Display")

    # ══════════════════════════════════════════════════════════════
    # SLIDE 7: BUSINESS MODEL  — white bg, 3 tier cards, right image
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)

    add_picture(slide, img("slide6_Picture 22.png"),
                Inches(9.6), Inches(0), Inches(3.73), Inches(7.5))

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5.0), Inches(0.4),
                 "BUSINESS MODEL", font_size=11, font_color=CORAL, font_name="DM Sans")

    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(8.5), Inches(0.7),
                 "Freemium with a clear upgrade path.",
                 font_size=34, font_color=BLACK, font_name="Playfair Display")

    coral_rule(slide, left=Inches(0.8), top=Inches(2.2))

    # Tier cards — FREE (white outline), PLUS (dark), PRO (coral)
    tiers = [
        ("FREE",     "$0",          [
            "3 AI checks / day",
            "3 follow-ups per check",
            "7-day history",
            "Ad-supported",
        ], "outline", BLACK, CHARCOAL),
        ("PLUS",     "$4.99/mo",    [
            "Unlimited AI checks",
            "5 follow-ups per check",
            "Full outfit history",
            "No ads",
            "Community feedback",
        ], BLACK, WHITE, GRAY),
        ("PRO",      "$14.99/mo",   [
            "Everything in Plus",
            "10 follow-ups per check",
            "5 expert reviews / month",
            "Event planning mode",
            "Style analytics & DNA",
        ], CORAL, WHITE, WHITE),
    ]

    for i, (name, price, features, bg, name_color, feat_color) in enumerate(tiers):
        x = Inches(0.8 + i * 3.0)
        if bg == "outline":
            add_outline_card(slide, x, Inches(2.6), Inches(2.6), Inches(4.2))
        else:
            add_shape(slide, x, Inches(2.6), Inches(2.6), Inches(4.2), bg)

        add_text_box(slide, x + Inches(0.2), Inches(2.8), Inches(2.2), Inches(0.3),
                     name, font_size=11, font_color=name_color, font_name="DM Sans")
        add_text_box(slide, x + Inches(0.2), Inches(3.1), Inches(2.2), Inches(0.6),
                     price, font_size=36, font_color=name_color, font_name="Playfair Display")
        card_divider(slide, x + Inches(0.2), Inches(3.8), Inches(2.2),
                     color=GRAY if bg != "outline" else DIVIDER)
        add_multiline_text(slide, x + Inches(0.2), Inches(4.0), Inches(2.2), Inches(2.5), [
            {"text": f, "size": 14, "color": feat_color, "space_after": 5} for f in features
        ])

    gray_rule(slide, left=Inches(0.8), top=Inches(6.85))
    add_text_box(slide, Inches(0.8), Inches(6.95), Inches(8.5), Inches(0.4),
                 "Layer 2 \u2014 Fashion Intelligence (B2B):  "
                 "Trend reports  \u2022  Trend API  \u2022  White-label SDK  \u2022  Custom research  "
                 "\u2014  highest-margin revenue at scale",
                 font_size=13, font_color=CORAL, font_name="DM Sans")

    # ══════════════════════════════════════════════════════════════
    # SLIDE 7: UNIT ECONOMICS  — white bg, 4 dark cards
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5.0), Inches(0.4),
                 "UNIT ECONOMICS", font_size=11, font_color=CORAL, font_name="DM Sans")

    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11.0), Inches(0.7),
                 "AI costs drop. Our margins improve. Every quarter.",
                 font_size=34, font_color=BLACK, font_name="Playfair Display")

    coral_rule(slide, left=Inches(0.8), top=Inches(2.2))

    metrics = [
        ("~$0.003",     "Cost per\nAI check",           "Gemini Vision API\ncost per analysis"),
        ("~$0.50",      "Blended monthly\ncost per user","Including infra,\nstorage, AI calls"),
        ("$4.99\u2013$14.99", "Monthly revenue\nper paid user", "Strong gross margins\non subscriptions"),
        ("10\u201330x",       "Target\nLTV:CAC",              "Organic-first\nacquisition"),
    ]
    for i, (value, label, detail) in enumerate(metrics):
        x = Inches(0.8 + i * 3.1)
        add_shape(slide, x, Inches(2.6), Inches(2.7), Inches(2.8), BLACK)
        add_text_box(slide, x + Inches(0.2), Inches(2.9), Inches(2.3), Inches(0.6),
                     value, font_size=32, font_color=WHITE, font_name="Playfair Display")
        add_text_box(slide, x + Inches(0.2), Inches(3.6), Inches(2.3), Inches(0.6),
                     label, font_size=14, font_color=WHITE, font_name="DM Sans")
        add_text_box(slide, x + Inches(0.2), Inches(4.3), Inches(2.3), Inches(0.6),
                     detail, font_size=12, font_color=GRAY, font_name="DM Sans")

    gray_rule(slide, left=Inches(0.8), top=Inches(5.8))
    add_text_box(slide, Inches(0.8), Inches(6.0), Inches(11.0), Inches(0.8),
                 "AI API costs have dropped 90% in 18 months and keep falling. "
                 "Our COGS improves automatically \u2014 the opposite of most consumer businesses.",
                 font_size=17, font_color=CHARCOAL, font_name="Playfair Display")

    # ══════════════════════════════════════════════════════════════
    # SLIDE 8: GROWTH  — white bg, left loop + dark GTM card, right image
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)

    add_picture(slide, img("slide8_Picture 8.jpg"),
                Inches(10.1), Inches(0), Inches(3.23), Inches(7.5))

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5.0), Inches(0.4),
                 "GROWTH", font_size=11, font_color=CORAL, font_name="DM Sans")

    add_text_box(slide, Inches(0.8), Inches(0.8), Inches(7.0), Inches(0.7),
                 "\u201cOr This?\u201d is inherently shareable.",
                 font_size=34, font_color=BLACK, font_name="Playfair Display")

    add_text_box(slide, Inches(0.8), Inches(2.0), Inches(7.0), Inches(0.4),
                 "Built-in virality. Every outfit is a potential share.",
                 font_size=16, font_color=CHARCOAL, font_name="DM Sans")

    coral_rule(slide, left=Inches(0.8), top=Inches(2.6))

    add_multiline_text(slide, Inches(0.8), Inches(2.9), Inches(4.5), Inches(3.0), [
        {"text": "The core loop:", "size": 17, "color": BLACK, "space_after": 8},
        {"text": "1.  User gets feedback \u2192 tells friends",  "size": 16, "color": CHARCOAL, "space_after": 6},
        {"text": "2.  A/B comparisons shared for votes",         "size": 16, "color": CHARCOAL, "space_after": 6},
        {"text": "3.  Friends download to participate",          "size": 16, "color": CHARCOAL, "space_after": 6},
        {"text": "4.  New users try AI \u2192 get hooked",       "size": 16, "color": CHARCOAL, "space_after": 6},
        {"text": "5.  Every outfit = potential share",           "size": 16, "color": CHARCOAL, "space_after": 6},
    ])

    # Dark GTM card
    add_shape(slide, Inches(5.8), Inches(2.9), Inches(4.0), Inches(3.8), BLACK)
    add_multiline_text(slide, Inches(6.1), Inches(3.1), Inches(3.5), Inches(3.4), [
        {"text": "Go-to-market:",          "size": 15, "color": WHITE,  "space_after": 8},
        {"text": "\u2022  TikTok / Reels", "size": 14, "color": GRAY,   "space_after": 5},
        {"text": "\u2022  Reddit (2M+ members)", "size": 14, "color": GRAY, "space_after": 5},
        {"text": "\u2022  Micro-creator seeding", "size": 14, "color": GRAY, "space_after": 5},
        {"text": "\u2022  Word of mouth",  "size": 14, "color": GRAY,   "space_after": 10},
        {"text": "Launch:",                "size": 15, "color": WHITE,  "space_after": 8},
        {"text": "\u2022  Private beta \u2192 500", "size": 14, "color": GRAY, "space_after": 5},
        {"text": "\u2022  ProductHunt launch",     "size": 14, "color": GRAY, "space_after": 5},
    ])

    # ══════════════════════════════════════════════════════════════
    # SLIDE 9: COMPETITION  — white bg, 4 cards (3 white + 1 coral)
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5.0), Inches(0.4),
                 "COMPETITION", font_size=11, font_color=CORAL, font_name="DM Sans")

    add_multiline_text(slide, Inches(0.8), Inches(1.2), Inches(11.0), Inches(0.8), [
        {"text": "We own the \u201cmoment of decision.\u201d", "size": 34, "color": BLACK, "font_name": "Playfair Display"},
        {"text": "Nobody else does.",                           "size": 34, "color": BLACK, "font_name": "Playfair Display"},
    ], line_spacing=1.15)

    coral_rule(slide, left=Inches(0.8), top=Inches(2.5))

    competitors = [
        ("Stitch Fix",  "Subscription boxes.",        "Days, not seconds.\nExpensive.",          "outline",  BLACK, CHARCOAL),
        ("Pinterest",   "Inspiration.",               "Not feedback.\nNot private.",             "outline",  BLACK, CHARCOAL),
        ("ChatGPT",     "General AI.",                "No fashion expertise.\nNo specialized UX.", "outline", BLACK, CHARCOAL),
        ("Or This?",    "Instant AI feedback\nat the moment of decision.", "Fast. Private.\nExpert. Always on.", CORAL, WHITE, WHITE),
    ]
    for i, (name, what, diff, bg, name_clr, body_clr) in enumerate(competitors):
        x = Inches(0.8 + i * 3.1)
        if bg == "outline":
            add_outline_card(slide, x, Inches(2.9), Inches(2.7), Inches(2.4))
        else:
            add_shape(slide, x, Inches(2.9), Inches(2.7), Inches(2.4), bg)
        add_text_box(slide, x + Inches(0.2), Inches(3.1), Inches(2.3), Inches(0.4),
                     name, font_size=17, font_color=name_clr, font_name="DM Sans")
        add_text_box(slide, x + Inches(0.2), Inches(3.6), Inches(2.3), Inches(0.5),
                     what, font_size=13, font_color=body_clr, font_name="DM Sans")
        add_text_box(slide, x + Inches(0.2), Inches(4.2), Inches(2.3), Inches(0.7),
                     diff, font_size=12, font_color=body_clr, font_name="DM Sans")

    gray_rule(slide, left=Inches(0.8), top=Inches(5.6))
    add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11.0), Inches(0.5),
                 "Our moat: The world\u2019s only real-time consumer fashion dataset \u2014 "
                 "every outfit analyzed generates structured intelligence that compounds in value. "
                 "Style DNA profiles that deepen over time. Community network effects.",
                 font_size=15, font_color=CHARCOAL, font_name="DM Sans")

    # ══════════════════════════════════════════════════════════════
    # SLIDE 10: TRACTION  — white bg, 4 dark status cards + targets
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5.0), Inches(0.4),
                 "TRACTION", font_size=11, font_color=CORAL, font_name="DM Sans")

    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11.0), Inches(0.7),
                 "Complete product built. Ready for users.",
                 font_size=34, font_color=BLACK, font_name="Playfair Display")

    coral_rule(slide, left=Inches(0.8), top=Inches(2.2))

    status_items = [
        ("BUILT",    "Full-stack app",    "37 screens\n16 API domains\n54+ database tables"),
        ("DEPLOYED", "Backend live",      "Railway hosting\nPostgreSQL\nAll services up"),
        ("READY",    "App store",         "iOS + Android\nbuilds prepared"),
        ("ACTIVE",   "AI pipeline",       "7-stage training\nPrompt v3.0"),
    ]
    for i, (badge, title, desc) in enumerate(status_items):
        x = Inches(0.8 + i * 3.1)
        add_shape(slide, x, Inches(2.5), Inches(2.7), Inches(2.2), BLACK)
        add_text_box(slide, x + Inches(0.2), Inches(2.7), Inches(2.3), Inches(0.3),
                     badge, font_size=11, font_color=CORAL, font_name="DM Sans")
        add_text_box(slide, x + Inches(0.2), Inches(3.0), Inches(2.3), Inches(0.3),
                     title, font_size=17, font_color=WHITE, font_name="DM Sans")
        add_text_box(slide, x + Inches(0.2), Inches(3.4), Inches(2.3), Inches(1.0),
                     desc, font_size=13, font_color=GRAY, font_name="DM Sans")

    gray_rule(slide, left=Inches(0.8), top=Inches(5.0))

    add_text_box(slide, Inches(0.8), Inches(5.2), Inches(5.0), Inches(0.3),
                 "MONTH 6 TARGETS", font_size=11, font_color=CORAL, font_name="DM Sans")

    targets = [
        ("50K",     "MAU"),
        ("10K",     "Daily checks"),
        ("30%",     "D7 retention"),
        ("4.0 / 5", "Helpfulness"),
    ]
    for i, (val, label) in enumerate(targets):
        x = Inches(0.8 + i * 3.1)
        add_text_box(slide, x, Inches(5.5), Inches(2.5), Inches(0.5),
                     val, font_size=32, font_color=BLACK, font_name="Playfair Display")
        add_text_box(slide, x, Inches(6.0), Inches(2.5), Inches(0.3),
                     label, font_size=13, font_color=CHARCOAL, font_name="DM Sans")

    # ══════════════════════════════════════════════════════════════
    # SLIDE 11: ROADMAP  — white bg, 4 phase cards
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5.0), Inches(0.4),
                 "ROADMAP", font_size=11, font_color=CORAL, font_name="DM Sans")

    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11.0), Inches(0.7),
                 "Nail the core. Then expand.",
                 font_size=34, font_color=BLACK, font_name="Playfair Display")

    coral_rule(slide, left=Inches(0.8), top=Inches(2.2))

    phases = [
        ("NOW",     "Core Loop",  CORAL,    WHITE, WHITE, GRAY, [
            "App store launch",
            "AI refinement",
            "Beta (500 users)",
            "Organic seeding",
        ]),
        ("Q3 2026", "Community",  BLACK,    WHITE, WHITE, GRAY, [
            "Community feedback",
            "\u201cOr This?\u201d A/B",
            "Give-to-get flywheel",
            "Gamification",
        ]),
        ("Q4 2026", "Monetize",   BLACK,    WHITE, WHITE, GRAY, [
            "Plus & Pro tiers",
            "Expert marketplace",
            "A/B test pricing",
            "Revenue optimization",
        ]),
        ("2027",    "Scale",      WHITE,    BLACK, BLACK, CHARCOAL, [
            "B2B trend API launch",
            "Fashion intelligence",
            "Shopping integration",
            "International",
        ]),
    ]

    for i, (when, title, bg, when_clr, title_clr, item_clr, items) in enumerate(phases):
        x = Inches(0.8 + i * 3.1)
        if bg == WHITE:
            add_outline_card(slide, x, Inches(2.5), Inches(2.7), Inches(4.5))
        else:
            add_shape(slide, x, Inches(2.5), Inches(2.7), Inches(4.5), bg)
        add_text_box(slide, x + Inches(0.2), Inches(2.7), Inches(2.3), Inches(0.3),
                     when, font_size=11, font_color=when_clr, font_name="DM Sans")
        add_text_box(slide, x + Inches(0.2), Inches(3.0), Inches(2.3), Inches(0.4),
                     title, font_size=22, font_color=title_clr, font_name="Playfair Display")
        add_multiline_text(slide, x + Inches(0.2), Inches(3.6), Inches(2.3), Inches(3.0), [
            {"text": f"\u2022  {item}", "size": 14, "color": item_clr, "space_after": 5}
            for item in items
        ])

    # ══════════════════════════════════════════════════════════════
    # SLIDE 12: THE TEAM  — white bg, left founder photo + right bio + top-right logos
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5.0), Inches(0.4),
                 "THE TEAM", font_size=11, font_color=CORAL, font_name="DM Sans")

    add_text_box(slide, Inches(0.3), Inches(1.18), Inches(11.0), Inches(0.64),
                 "Built by someone who understood the problem.",
                 font_size=34, font_color=BLACK, font_name="Playfair Display")

    coral_rule(slide, left=Inches(0.8), top=Inches(2.2))

    # Left — founder photo placeholder (dark rect same size as reference)
    add_picture(slide, img("slide12_Picture 26.jpg"),
                Inches(0.3), Inches(2.43), Inches(3.38), Inches(5.07),
                placeholder=BLACK)

    # Right — name + bio
    add_text_box(slide, Inches(4.2), Inches(2.47), Inches(5.2), Inches(0.4),
                 "Brandon Davis", font_size=18, font_color=BLACK, font_name="DM Sans",
                 alignment=PP_ALIGN.LEFT)
    add_text_box(slide, Inches(4.2), Inches(2.87), Inches(5.2), Inches(0.3),
                 "CEO / Product", font_size=13, font_color=CORAL, font_name="DM Sans")
    add_text_box(slide, Inches(4.2), Inches(3.27), Inches(5.2), Inches(0.3),
                 "Vassar College  \u2022  Kellogg MBA (Northwestern)",
                 font_size=13, font_color=CHARCOAL, font_name="DM Sans")

    add_multiline_text(slide, Inches(4.2), Inches(3.9), Inches(5.5), Inches(3.5), [
        {"text": "Kellogg MBA (Northwestern). Led program management at Apple "
                 "(Health Technologies), Verily/Alphabet (Project Baseline), and Komodo Health. "
                 "8+ years AI product management. $15M+ revenue impact.",
         "size": 14, "color": CHARCOAL, "space_after": 10},
        {"text": "Former signed model (Rae Agency) and SAG-AFTRA actor. "
                 "Lived the problem \u2014 understands fashion from the inside.",
         "size": 14, "color": CHARCOAL, "space_after": 10},
        {"text": "Technical founder: built this entire product \u2014 "
                 "React Native, TypeScript, PostgreSQL.",
         "size": 14, "color": CHARCOAL},
    ], line_spacing=1.4)

    # Top-right company logos
    logo_imgs = [
        ("slide12_Picture 2.png",  Inches(9.35),  Inches(0.13), Inches(1.28), Inches(1.28)),
        ("slide12_Picture 4.jpg",  Inches(10.63), Inches(0.02), Inches(2.67), Inches(1.27)),
        ("slide12_Picture 8.png",  Inches(10.35), Inches(1.29), Inches(2.67), Inches(0.94)),
        ("slide12_Picture 10.png", Inches(10.67), Inches(2.4),  Inches(2.04), Inches(1.14)),
        ("slide12_Picture 12.png", Inches(10.23), Inches(3.94), Inches(2.92), Inches(0.81)),
        ("slide12_Picture 14.png", Inches(10.81), Inches(4.92), Inches(2.34), Inches(2.34)),
    ]
    for fname, l, t, w, h in logo_imgs:
        path = img(fname)
        if path:
            slide.shapes.add_picture(path, l, t, w, h)

    # ══════════════════════════════════════════════════════════════
    # SLIDE 13: THE ASK  — dark bg, 4 white fund cards, right image
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BLACK)

    add_picture(slide, img("slide13_Picture 4.jpg"),
                Inches(10.0), Inches(0), Inches(3.33), Inches(7.5))

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5.0), Inches(0.4),
                 "THE ASK", font_size=11, font_color=GRAY, font_name="DM Sans")

    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(8.0), Inches(0.7),
                 "Raising $1.5M seed to launch and scale.",
                 font_size=36, font_color=WHITE, font_name="Playfair Display")

    coral_rule(slide, left=Inches(0.8), top=Inches(2.2))

    funds = [
        ("40%", "Engineering",  "Scale infra.\nHire 2 engineers.",     Inches(0.8)),
        ("25%", "Growth",       "Creator partnerships.\nLaunch.",       Inches(3.1)),
        ("20%", "AI / Data",    "Training pipeline.\nStyle DNA.",       Inches(5.4)),
        ("15%", "Operations",   "Team. Legal.\n18-mo runway.",          Inches(7.7)),
    ]
    for pct, label, desc, x in funds:
        add_outline_card(slide, x, Inches(2.6), Inches(2.0), Inches(2.6))
        add_text_box(slide, x + Inches(0.15), Inches(2.8), Inches(1.7), Inches(0.5),
                     pct, font_size=32, font_color=CORAL, font_name="Playfair Display")
        add_text_box(slide, x + Inches(0.15), Inches(3.3), Inches(1.7), Inches(0.3),
                     label, font_size=14, font_color=BLACK, font_name="DM Sans")
        card_divider(slide, x + Inches(0.15), Inches(3.6), Inches(1.7))
        add_text_box(slide, x + Inches(0.15), Inches(3.8), Inches(1.7), Inches(1.0),
                     desc, font_size=12, font_color=CHARCOAL, font_name="DM Sans")

    gray_rule(slide, left=Inches(0.8), top=Inches(5.6), width=Inches(8.5))
    add_text_box(slide, Inches(0.8), Inches(5.8), Inches(8.0), Inches(0.3),
                 "18-MONTH MILESTONES", font_size=11, font_color=GRAY, font_name="DM Sans")
    add_text_box(slide, Inches(0.8), Inches(6.1), Inches(8.5), Inches(0.5),
                 "250K MAU  \u2022  75K daily checks  \u2022  $1M+ ARR  \u2022  Series A metrics",
                 font_size=17, font_color=WHITE, font_name="DM Sans")

    # ══════════════════════════════════════════════════════════════
    # SLIDE 14: CLOSING  — dark bg, left half image, right logo + contact
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BLACK)

    add_picture(slide, img("slide14_Picture 1.jpg"),
                Inches(0), Inches(0), Inches(5.5), Inches(7.5))

    coral_rule(slide, left=Inches(6.5), top=Inches(2.2), width=Inches(2.0))

    add_text_box(slide, Inches(6.5), Inches(2.5), Inches(3.0), Inches(1.0),
                 "Or", font_size=72, font_color=WHITE, font_name="DM Sans")
    add_text_box(slide, Inches(8.0), Inches(2.47), Inches(4.0), Inches(1.0),
                 "This?", font_size=72, font_color=CORAL, font_name="Playfair Display")

    add_text_box(slide, Inches(6.5), Inches(3.7), Inches(6.0), Inches(0.5),
                 "Confidence in every choice.",
                 font_size=24, font_color=WHITE, font_name="Playfair Display")

    add_text_box(slide, Inches(6.5), Inches(5.0), Inches(5.0), Inches(0.4),
                 "bradavis2011@gmail.com  \u2022  orthis.app",
                 font_size=15, font_color=GRAY, font_name="DM Sans")

    # ── Save ──
    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                               "OrThis_Seed_Pitch_Deck.pptx")
    prs.save(output_path)
    print(f"Pitch deck saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    build_deck()
